import streamlit as st
st.set_page_config(page_title="AICOSS Chatbot (AICOSS 챗봇)", layout="wide")

import os
import json
import time
import numpy as np
import faiss
import re
from sentence_transformers import SentenceTransformer
from transformers import MarianMTModel, MarianTokenizer
import ollama
from PIL import Image

# ========================================
# [추가] 파일 파싱을 위한 라이브러리
# ========================================
import PyPDF2
from docx import Document
from pptx import Presentation
import tempfile

import warnings
warnings.filterwarnings("ignore")

# =====================================================
# 1. 기본 설정: 경로 및 학교 별칭, 기타 유틸리티 함수들
# =====================================================
base_dir = r"E:\deep_learning\THU\Capstone\7univ\02057\Capstone"  # 본인 환경에 맞게 수정
combined_data_path = os.path.join(base_dir, "combined_data.json")
faiss_index_path = os.path.join(base_dir, "faiss_index.bin")

# 학교별 별칭(정규표현식 패턴은 모두 raw 문자열로)
school_aliases = {
    "jju": [
        r"전\s*주\s*대학교", r"전\s*주\s*대", r"전\s*대", r"전\s*주\s*대학",
        r"jeonju\s*uni", r"jeonju\s*university", r"JJU"
    ],
    "jnu": [
        r"전\s*남\s*대학교", r"전\s*남\s*대", r"전\s*남\s*대학", r"전남대",
        r"jeonnam\s*uni", r"jeonnam\s*university", r"jeonnam\s*national\s*university", r"JNU"
    ],
    "knu": [
        r"경\s*북\s*대학교", r"경\s*북\s*대", r"경\s*북\s*대학", r"경북대",
        r"kyungpook\s*uni", r"kyungpook\s*university", r"kyungpook\s*national\s*university", r"KNU"
    ],
    "skku": [
        r"성\s*균\s*관\s*대학교", r"성\s*균\s*관\s*대", r"성\s*대", r"성균관대",
        r"sungkyunkwan\s*uni", r"sungkyunkwan\s*university", r"SKKU"
    ],
    "tec": [
        r"서\s*울\s*과\s*학\s*기\s*술\s*대학교", r"서\s*울\s*과\s*학\s*기\s*술\s*대",
        r"과\s*학\s*기\s*술\s*대학교", r"과\s*학\s*기\s*술\s*대",
        r"과\s*기\s*대학교", r"과\s*기\s*대", r"서울과기대", r"SeoulTech",
        r"Seoul\s*National\s*University\s*of\s*Science\s*and\s*Technology",
        r"Seoul\s*National\s*University", r"TEC"
    ],
    "uos": [
        r"서\s*울\s*시\s*립\s*대학교", r"서\s*울\s*시\s*립\s*대", r"시립대",
        r"uni\s*of\s*Seoul", r"University\s*of\s*Seoul", r"UOS"
    ],
    "yj": [
        r"영\s*진\s*전\s*문\s*대학교", r"영\s*진\s*전\s*문\s*대", r"영진대", r"영진전문대",
        r"Yeungjin\s*University", r"Yeungjin\s*uni", r"YJU"
    ]
}

def normalize_school_names(text: str) -> str:
    """학교 이름의 다양한 변형을 표준 키로 치환합니다."""
    for standard, aliases in school_aliases.items():
        for alias in aliases:
            text = re.sub(alias, standard, text, flags=re.IGNORECASE)
    return text

def exclude_school_names_during_translation(text: str):
    """번역 전 학교 이름을 보호하여 번역되지 않도록 합니다."""
    protected_names = {}
    for standard, aliases in school_aliases.items():
        for alias in aliases:
            matches = re.findall(alias, text, flags=re.IGNORECASE)
            for match in matches:
                if match not in protected_names.values():
                    placeholder = f"[[{standard}]]"
                    protected_names[placeholder] = match
                    text = text.replace(match, placeholder)
    return text, protected_names

def restore_protected_names(text: str, protected_names: dict) -> str:
    """번역 후 보호된 학교 이름을 원래대로 복원합니다."""
    for placeholder, original in protected_names.items():
        text = text.replace(placeholder, original)
    return text

def get_school_key(school_name: str) -> str:
    """
    입력된 학교 이름 문자열이 school_aliases에 정의된 패턴과 매칭되면,
    해당 표준 키(예: "jju")를 반환합니다. 매칭되지 않으면 원래 문자열을 반환합니다.
    """
    for key, alias_list in school_aliases.items():
        if school_name.lower() == key.lower():
            return key
        for alias in alias_list:
            if re.fullmatch(alias, school_name, flags=re.IGNORECASE):
                return key
    return school_name

# =====================================================
# 2. 모델 및 데이터 로딩 (st.cache_resource 사용)
# =====================================================
@st.cache_resource
def load_all_entries():
    with open(combined_data_path, 'r', encoding='utf-8') as f:
        combined_data = json.load(f)
    all_entries = combined_data.get("data_entries", []) + combined_data.get("image_entries", [])
    return all_entries

@st.cache_resource
def load_faiss_index():
    index = faiss.read_index(faiss_index_path)
    st.write(f"Loaded FAISS index, total vectors: {index.ntotal} (불러온 FAISS 인덱스, 총 벡터 수: {index.ntotal})")
    return index

@st.cache_resource
def load_translation_model():
    model_name = "Helsinki-NLP/opus-mt-ko-en"
    tokenizer = MarianTokenizer.from_pretrained(model_name)
    model = MarianMTModel.from_pretrained(model_name)
    return tokenizer, model

@st.cache_resource
def load_embedding_model():
    model = SentenceTransformer('all-MiniLM-L6-v2')
    return model

# 전역 변수 로딩
all_entries = load_all_entries()
faiss_index = load_faiss_index()
trans_tokenizer, trans_model = load_translation_model()
embedding_model = load_embedding_model()

# =====================================================
# 3. 번역 및 임베딩 관련 함수
# =====================================================
def translate_ko_to_en(text: str, tokenizer, model) -> str:
    """
    한국어 텍스트를 영어로 번역하되, 학교 이름(정규식에 맞는 부분)은 번역하지 않고
    school_aliases에 정의된 키 값(예: "jju")으로 변경합니다.
    """
    patterns = []
    for key, alias_list in school_aliases.items():
        patterns.append(re.escape(key))
        patterns.extend(alias_list)
    combined_pattern = "(" + "|".join(patterns) + ")"
    pattern_re = re.compile(combined_pattern, flags=re.IGNORECASE)

    parts = re.split(pattern_re, text)
    translated_parts = []
    for i, part in enumerate(parts):
        if i % 2 == 1:
            # 학교 이름(혹은 별칭) 부분이면 키로 normalize
            normalized_key = get_school_key(part)
            translated_parts.append(normalized_key)
        else:
            if part.strip() == "":
                translated_parts.append(part)
            else:
                inputs = tokenizer(part, return_tensors="pt", padding=True, truncation=True)
                translated = model.generate(**inputs)
                translated_text = tokenizer.decode(translated[0], skip_special_tokens=True)
                translated_parts.append(translated_text)
    return "".join(translated_parts)

def get_embedding(text: str):
    emb = embedding_model.encode(text, convert_to_numpy=True)
    return emb.astype('float32')

def detect_language(text: str) -> str:
    hangul_count = sum(1 for ch in text if '가' <= ch <= '힣')
    english_count = sum(1 for ch in text if ch.lower() in 'abcdefghijklmnopqrstuvwxyz')
    return 'ko' if hangul_count >= english_count else 'en'

# =====================================================
# 4. 검색 및 Q&A 관련 함수
# =====================================================
def get_university_from_question(question: str):
    """
    질문에서 학교 이름(대학 명칭)을 찾습니다.
    """
    for standard, aliases in school_aliases.items():
        for alias in aliases:
            if re.search(alias, question, flags=re.IGNORECASE):
                return standard
    return None

def retrieve_best_matches(query: str, lang: str, top_k: int = 3):
    """
    내부 DB에서 검색해올 때 사용.
    """
    university_key = get_university_from_question(query)
    if university_key is None:
        return None, [], "⚠ Unable to find university name in question (⚠ 질문에서 대학 이름을 찾을 수 없습니다.)"
    
    filtered_entries = [entry for entry in all_entries if entry.get("university") == university_key]
    if not filtered_entries:
        return None, [], f"⚠ No data available for {university_key} (⚠ {university_key}에 대한 데이터가 없습니다.)"
    
    index_to_content_kr = {entry["index"]: entry["content_kr"] for entry in filtered_entries if "content_kr" in entry}
    index_to_content_en = {entry["index"]: entry["content"] for entry in filtered_entries if "content" in entry}
    index_to_summary = {entry["index"]: entry.get("summary", "(요약 없음)") for entry in filtered_entries}

    summary_vectors = np.array(
        [embedding_model.encode(index_to_summary[idx]) for idx in index_to_summary.keys()],
        dtype='float32'
    )
    # FAISS 서브 인덱스
    faiss_subset = faiss.IndexFlatL2(summary_vectors.shape[1])
    faiss_subset.add(summary_vectors)
    valid_indices = list(index_to_summary.keys())

    if lang == 'ko':
        translated_query = translate_ko_to_en(query, trans_tokenizer, trans_model)
        search_query = translated_query
    else:
        translated_query = None
        search_query = normalize_school_names(query)

    query_emb = embedding_model.encode(search_query)
    query_emb = np.expand_dims(query_emb, axis=0)

    distances, subset_indices = faiss_subset.search(query_emb, top_k)
    retrieved_contexts = []
    retrieved_indices = []
    for subset_idx in subset_indices[0]:
        real_idx = valid_indices[subset_idx]
        retrieved_indices.append(real_idx)
        content = index_to_content_kr.get(real_idx) if lang == 'ko' else index_to_content_en.get(real_idx)
        summary = index_to_summary.get(real_idx)
        if content is not None:
            # 여기서는 content + summary를 묶어서 "reference_info" 로 사용
            retrieved_contexts.append(f"{content}\n\nexplanation:\n{summary}")
    full_reference = "\n\n".join(retrieved_contexts)
    return translated_query, retrieved_indices, full_reference

# =====================================================
# === 새롭게 분리된 두 함수: generate_answer_ephemeral / generate_answer_db
# =====================================================
def generate_answer_ephemeral(user_msg: str, conversation_history: str, reference_info: str, lang: str) -> str:
    """
    [파일 기반 Q&A]를 위한 함수.
    업로드된 파일에서 가져온 참고 정보를 활용해 답변.
    """
    OLLAMA_MODEL = "qwen2.5:1.5b"

    if lang == 'ko':
        prompt = f"""
사용자의 질문에 대화 내역과 참고 정보를 참고하여 한국어로 답변하세요. (Answer in KOREAN)
정확한 정보인 "참고 정보" 안에 있는 내용을 사용해서 답변하세요.
거짓말은 할 수 없으며 친절하게 대답해야 해요.

대학교명을 사용하는 경우, "학교명"에 있는 학교명만 사용하세요.

학교명: 전주대학교, 전남대학교, 경북대학교, 성균관대학교, 서울과학기술대학교, 서울시립대학교, 영진전문대학교

[대화 내역]:
{conversation_history}

[참고 정보]:
{reference_info}

질문:
{user_msg}
"""
    else:
        prompt = f"""
[Source: User uploaded file]
Answer the user's question using the provided conversation history and reference information.
Use Reference Information which is trustable to answer.
You can't lie and should answer politely.

When you talk about university, only use University Name in "University Name"

University Name: 
Jeonju University, Jeonnam National University, Kyungpook National University, Sungkyunkwan University, Seoul National University of Science and Technology, University of Seoul, Yeungjin University


Conversation History:
{conversation_history}

Reference Information:
{reference_info}

Question:
{user_msg}
"""

    response = ollama.chat(model=OLLAMA_MODEL, messages=[{"role": "user", "content": prompt}])
    return response["message"]["content"]

def generate_answer_db(user_msg: str, conversation_history: str, reference_info: str, lang: str) -> str:
    """
    [내부 DB 기반 Q&A]를 위한 함수.
    내부 DB에서 가져온 참고 정보를 활용해 답변.
    """
    OLLAMA_MODEL = "qwen2.5:1.5b"

    if lang == 'ko':
        prompt = f"""
사용자의 질문에 대화 내역과 참고 정보를 참고하여 한국어로 답변하세요. (Answer in KOREAN)
"참고 정보" 안에 있는 내용을 중점적으로 활용하세요.
거짓말은 할 수 없으며 친절하게 대답해야 해요.

대학교명을 사용하는 경우, "학교명"에 있는 학교명만 사용하세요.

학교명: 전주대학교, 전남대학교, 경북대학교, 성균관대학교, 서울과학기술대학교, 서울시립대학교, 영진전문대학교


[대화 내역]:
{conversation_history}

[참고 정보]:
{reference_info}

질문:
{user_msg}
"""
    else:
        prompt = f"""
[Source: Internal DB]
Answer the user's question using the provided conversation history and reference information.

When you talk about university, only use University Name in "University Name"

University Name: 
Jeonju University, Jeonnam National University, Kyungpook National University, Sungkyunkwan University, Seoul National University of Science and Technology, University of Seoul, Yeungjin University

Conversation History:
{conversation_history}

Reference Information:
{reference_info}

Question:
{user_msg}
"""

    response = ollama.chat(model=OLLAMA_MODEL, messages=[{"role": "user", "content": prompt}])
    return response["message"]["content"]

# =====================================================
# 5. 이미지 관련 함수 (기존 로직 유지)
# =====================================================
def retrieve_best_image(query: str, lang: str):
    """
    해당 대학에 속하는 image_entries 중에서, 이미지 summary를 임베딩하여
    코사인 유사도가 가장 높은 이미지 1장을 반환 (없으면 None)
    """
    university_key = get_university_from_question(query)
    if university_key is None:
        return None, None
    
    filtered_images = [entry for entry in all_entries 
                       if entry.get("university") == university_key 
                       and entry.get("data_type") == "image"]
    if not filtered_images:
        return None, None
    
    index_to_summary = {}
    for entry in filtered_images:
        idx = entry["index"]
        index_to_summary[idx] = entry.get("summary", "(no summary)")

    if lang == 'ko':
        translated_query = translate_ko_to_en(query, trans_tokenizer, trans_model)
        search_query = translated_query
    else:
        search_query = normalize_school_names(query)

    query_emb = embedding_model.encode(search_query)
    query_emb = np.expand_dims(query_emb, axis=0)

    image_vectors = np.array(
        [embedding_model.encode(index_to_summary[i]) for i in index_to_summary.keys()],
        dtype='float32'
    )
    faiss_subset = faiss.IndexFlatL2(image_vectors.shape[1])
    faiss_subset.add(image_vectors)
    valid_indices = list(index_to_summary.keys())

    distances, subset_indices = faiss_subset.search(query_emb, 1)
    best_idx = subset_indices[0][0]
    real_idx = valid_indices[best_idx]
    best_image_entry = None
    for ent in filtered_images:
        if ent["index"] == real_idx:
            best_image_entry = ent
            break

    if best_image_entry is None:
        return None, None

    best_image_summary = best_image_entry.get("summary", "(no summary)")
    return best_image_entry, best_image_summary

def ask_llm_if_use_image(user_msg: str, text_answer: str, image_summary: str) -> str:
    """
    LLM에게 이미지가 유저의 질문 답변에 꼭 필요한지 묻습니다.
    응답에 "yes"가 포함되면 이미지를 사용합니다.
    """
    OLLAMA_MODEL = "qwen2.5:1.5b"
    prompt = f"""
You must determine whether the given [Image Summary] has similar content or context with the [User Question].

[User Question]
{user_msg}

[Image Summary]
{image_summary}

Answer strictly with "yes" if the [Image Summary] is related to User Question.
Respond with only "yes" or "no".
"""
    response = ollama.chat(model=OLLAMA_MODEL, messages=[{"role": "user", "content": prompt}])
    content = response["message"]["content"].strip()
    return content

def get_relative_image_path(absolute_path: str) -> str:
    """
    절대경로를 Capstone 이후 하위 폴더부터의 상대경로로 변환합니다.
    """
    splitted = absolute_path.split('Capstone')
    if len(splitted) < 2:
        return os.path.basename(absolute_path)
    rel_part = splitted[-1].replace('\\', '/').lstrip('/')
    return rel_part

# =====================================================
# [추가 기능] PDF/Word/PPT 텍스트 추출 및 임베딩
# =====================================================
def is_english_paragraph(text: str, threshold: float = 0.7) -> bool:
    """
    문단에서 영어 알파벳(a-z, A-Z)이 차지하는 비율이 threshold 이상인지 판별.
    숫자나 공백, 특수문자는 제외하고 계산.
    """
    filtered = "".join(ch for ch in text if ch.isalpha())
    if len(filtered) == 0:
        return False
    english_chars = sum(1 for ch in filtered if ('a' <= ch.lower() <= 'z'))
    ratio = english_chars / len(filtered)
    return ratio >= threshold

def parse_pdf(file_content: bytes) -> str:
    """
    PDF에서 텍스트를 추출 (PyPDF2 사용).
    """
    text = ""
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp:
        tmp.write(file_content)
        tmp_path = tmp.name

    with open(tmp_path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"

    os.remove(tmp_path)
    return text

def parse_docx(file_content: bytes) -> str:
    """
    Word(docx)에서 텍스트를 추출 (python-docx 사용).
    """
    text = ""
    with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp:
        tmp.write(file_content)
        tmp_path = tmp.name
    
    doc = Document(tmp_path)
    for para in doc.paragraphs:
        text += para.text + "\n"
    os.remove(tmp_path)
    return text

def parse_pptx(file_content: bytes) -> str:
    """
    PPT(pptx)에서 텍스트 추출 (python-pptx 사용).
    """
    text = ""
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
        tmp.write(file_content)
        tmp_path = tmp.name

    prs = Presentation(tmp_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    os.remove(tmp_path)
    return text

# =====================================================
# [개선된] 문서 청킹 함수 (맥락기반, 문단+문장 혼합, 길이 제한)
# =====================================================
def chunk_text_smart(text: str, max_chars:int=1500, min_chars:int=700, overlap:int=200) -> list:
    """
    1) 빈 줄을 기준으로 문단을 나눔
    2) 각 문단이 너무 길면, 문장 단위로 분할
    3) 하나의 청크를 만들 때, 문단을 차례로 붙여가며 
       - 현재 청크 길이가 min_chars 이하일 때는 계속 붙임
       - max_chars를 넘기면 청크 확정
    4) 청크 간 overlap(문자 수)만큼 앞선 내용에서 겹쳐서 추가 (맥락 유지)
    """

    # 1) 빈 줄 기준으로 문단 나누기
    paragraphs = re.split(r'\n\s*\n', text.strip())
    paragraphs = [p.strip() for p in paragraphs if p.strip()]

    # 2) 만약 어떤 문단이 max_chars보다 훨씬 길면, 문장 단위 재분할
    expanded_paragraphs = []
    for para in paragraphs:
        if len(para) > max_chars:
            # 문장 단위로 쪼개서 max_chars 이하로 나눔
            sentences = re.split(r'(?<=[.?!])\s+', para)
            temp_buf = ""
            for sent in sentences:
                if len(temp_buf) + len(sent) <= max_chars:
                    temp_buf += (sent + " ")
                else:
                    expanded_paragraphs.append(temp_buf.strip())
                    temp_buf = sent + " "
            if temp_buf.strip():
                expanded_paragraphs.append(temp_buf.strip())
        else:
            expanded_paragraphs.append(para)

    # 3) 청크 생성
    chunks = []
    current_chunk = ""
    for para in expanded_paragraphs:
        if not current_chunk:
            # 비어 있으면 그냥 시작
            current_chunk = para
        else:
            # 이어붙였을 때 길이 확인
            if len(current_chunk) + len(para) + 1 <= max_chars:
                current_chunk += "\n" + para
            else:
                # 일단 현재 청크 확정
                chunks.append(current_chunk.strip())
                current_chunk = para

        # 만약 현재 청크가 최소 길이를 넘는 상태에서 max_chars 가까이 갔다면 바로 확정
        if len(current_chunk) > max_chars:
            chunks.append(current_chunk.strip())
            current_chunk = ""

    # 마지막 남은 청크
    if current_chunk.strip():
        chunks.append(current_chunk.strip())

    # 4) 오버랩 처리
    if overlap > 0:
        final_chunks = []
        for i, c in enumerate(chunks):
            if i == 0:
                final_chunks.append(c)
            else:
                # 이전 청크 뒤 overlap 만큼 잘라서 붙이기
                prev_chunk = final_chunks[-1]
                overlap_text = prev_chunk[-overlap:] if len(prev_chunk) > overlap else prev_chunk
                new_chunk = overlap_text + "\n" + c
                final_chunks.append(new_chunk)
        return final_chunks
    else:
        return chunks

def create_ephemeral_faiss_index(chunks: list) -> faiss.IndexFlatL2:
    """
    주어진 텍스트 청크들에 대해 임베딩을 계산하고 FAISS 인덱스를 만듦.
    """
    embeddings = [embedding_model.encode(c) for c in chunks]
    embeddings = np.array(embeddings, dtype='float32')
    index = faiss.IndexFlatL2(embeddings.shape[1])
    index.add(embeddings)
    return index

def ephemeral_retrieve_query(query: str, chunks: list, faiss_index: faiss.IndexFlatL2, top_k: int = 3) -> str:
    """
    업로드된 파일 청크에 대해서만 질문(query)을 검색하고 상위 top_k개를 반환.
    """
    query_emb = embedding_model.encode(query)
    query_emb = np.expand_dims(query_emb, axis=0)
    distances, indices = faiss_index.search(query_emb, top_k)

    retrieved_texts = []
    for idx in indices[0]:
        if 0 <= idx < len(chunks):
            retrieved_texts.append(chunks[idx])
    full_reference = "\n\n".join(retrieved_texts)
    return full_reference

# =====================================================
# 7. Streamlit UI 메인
# =====================================================
st.title("AICOSS Chatbot🤖")

# 대화 내역은 session_state에 저장
if "conversation_history" not in st.session_state:
    st.session_state.conversation_history = ""

# [추가] 파일 업로더 (PDF, Word, PPT) - 최대 1MB, 영어만 가능
uploaded_file = st.file_uploader(
    "Upload a PDF/Word/PPT file (Max:1MB, English only)",
    type=["pdf", "docx", "pptx"],
    help="Only English documents are supported. Max file size = 1MB."
)

with st.form("chat_form", clear_on_submit=True):
    user_input = st.text_input("Enter your question (질문을 입력하세요)")
    submitted = st.form_submit_button("Send (전송)")

# [추가] 업로드된 파일을 처리할 임시 변수
file_chunks = None
ephemeral_index = None

if uploaded_file:
    if uploaded_file.size > 1 * 1024 * 1024:
        st.warning("File size exceeds 1MB limit. The file will be ignored. (1MB 초과 파일은 무시됩니다.)")
    else:
        # 파일 파싱
        file_type = uploaded_file.name.split(".")[-1].lower()
        file_text = ""
        if file_type == "pdf":
            file_text = parse_pdf(uploaded_file.read())
        elif file_type == "docx":
            file_text = parse_docx(uploaded_file.read())
        elif file_type == "pptx":
            file_text = parse_pptx(uploaded_file.read())

        # 문단(줄바꿈 기준) 단위로 영어문단인지 확인하여 필터링
        paragraphs = file_text.split("\n")
        english_paragraphs = [p for p in paragraphs if is_english_paragraph(p.strip())]
        if not english_paragraphs:
            st.warning("No English paragraphs found or less than 70% English text. (영어 문단이 충분치 않습니다.)")
        else:
            # 유효한 문단만 모아 큰 텍스트로 만든 뒤, 개선된 청킹 함수 사용
            filtered_text = "\n\n".join(english_paragraphs)
            # 아래 파라미터들은 상황에 맞게 조절 가능
            file_chunks = chunk_text_smart(filtered_text, max_chars=1500, min_chars=700, overlap=200)
            ephemeral_index = create_ephemeral_faiss_index(file_chunks)

if submitted and user_input.strip():
    lang = detect_language(user_input)

    # [파일이 업로드된 경우] => 업로드 파일 기반 Q&A
    if file_chunks and ephemeral_index:
        with st.spinner("Searching your uploaded file..."):
            # 파일 질의 시, 영어 쿼리를 그대로 사용하거나 / 한국어면 영어로 번역
            if lang == 'ko':
                search_query = translate_ko_to_en(user_input, trans_tokenizer, trans_model)
            else:
                search_query = user_input  # 영어라면 그대로
            reference_info = ephemeral_retrieve_query(search_query, file_chunks, ephemeral_index, top_k=5)

        if not reference_info.strip():
            st.warning("⚠ Could not find relevant info in the uploaded file.")
        else:
            with st.spinner("Generating answer from your file..."):
                # 파일 기반 Q&A
                answer = generate_answer_ephemeral(
                    user_input,
                    st.session_state.conversation_history,
                    reference_info,
                    lang
                )

            # 스트리밍 형태 출력
            answer_placeholder = st.empty()
            streamed_answer = ""
            for char in answer:
                streamed_answer += char
                answer_placeholder.markdown(streamed_answer)
                time.sleep(0.02)

            # 사용된 참고 정보 **원본 그대로** 표시
            st.markdown("### Used Retrieved Information (사용된 참고 정보)")
            st.markdown(reference_info)

            # 대화 내역 업데이트
            st.session_state.conversation_history += f"User: {user_input}\nBot: {answer}\n\n"
            st.markdown("### Conversation History (대화 내역)")
            st.text_area("", st.session_state.conversation_history, height=200)

    else:
        # [파일이 없거나(크기 초과, 영어 아님, 등) -> DB 기반 Q&A]
        with st.spinner("Searching the internal DB..."):
            translated_query, indices, reference_info = retrieve_best_matches(user_input, lang, top_k=3)

        if not reference_info.strip():
            st.warning("⚠ Unable to find relevant information from DB.")
        else:
            with st.spinner("Generating answer from internal DB..."):
                answer = generate_answer_db(
                    user_input,
                    st.session_state.conversation_history,
                    reference_info,
                    lang
                )

            # 이미지 검색
            image_entry, image_summary = retrieve_best_image(user_input, lang)
            use_image = False
            relative_image_path = None
            if image_entry and image_summary:
                yes_no = ask_llm_if_use_image(user_input, answer, image_summary)
                if "yes" in yes_no.lower():
                    use_image = True
                    absolute_image_path = image_entry["file_path"]
                    relative_image_path = get_relative_image_path(absolute_image_path)

            # 이미지가 사용될 경우 표시
            if use_image and relative_image_path:
                image_absolute_path = os.path.join(base_dir, relative_image_path.replace('/', os.sep))
                try:
                    img = Image.open(image_absolute_path)
                    st.image(img, caption="Image (이미지)", use_column_width=True)
                except Exception as e:
                    st.error(f"Failed to display image: {e}")

            # 답변 스트리밍 출력
            answer_placeholder = st.empty()
            streamed_answer = ""
            for char in answer:
                streamed_answer += char
                answer_placeholder.markdown(streamed_answer)
                time.sleep(0.02)

            # 사용된 참고 정보 **원본 그대로** 표시
            st.markdown("### Used Retrieved Information (사용된 참고 정보)")
            st.markdown(reference_info)

            # 대화 내역 업데이트
            st.session_state.conversation_history += f"User: {user_input}\nBot: {answer}\n\n"
            st.markdown("### Conversation History (대화 내역)")
            st.text_area("", st.session_state.conversation_history, height=200)
